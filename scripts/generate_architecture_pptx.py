"""
Generate an architecture/code-structure PowerPoint for TESTTHEKEY.
Slides use monospace ASCII diagrams for clear code comprehension.

Run from project root:
    python scripts/generate_architecture_pptx.py
Output: TESTTHEKEY_Architecture.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_PANEL   = RGBColor(0x13, 0x25, 0x38)
C_CODE    = RGBColor(0x0A, 0x12, 0x1C)   # near-black for code blocks
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_GREEN   = RGBColor(0x2E, 0xCC, 0x71)
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)
C_YELLOW  = RGBColor(0xF3, 0x9C, 0x12)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xB0, 0xC4, 0xDE)
C_PURPLE  = RGBColor(0x9B, 0x59, 0xB6)

W, H = Inches(13.33), Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _bg(slide, color: RGBColor = C_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, bg=None, border=None, border_w=1.2):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.background() if not bg else (s.fill.solid(), setattr(s.fill.fore_color, 'rgb', bg))
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def _txt(slide, text, l, t, w, h,
         size=14, bold=False, color=C_WHITE,
         align=PP_ALIGN.LEFT, italic=False, font="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return txb


def _mono(slide, text, l, t, w, h, size=11, color=C_GREEN, bg=C_CODE, wrap=True):
    """Monospace code block with dark background."""
    _rect(slide, l, t, w, h, bg=bg)
    _txt(slide, text, l + Inches(0.12), t + Inches(0.1),
         w - Inches(0.2), h - Inches(0.15),
         size=size, color=color, font="Courier New", wrap=wrap)


def _header(slide, title, subtitle=""):
    _rect(slide, 0, 0, Inches(0.2), H, bg=C_ACCENT)
    _txt(slide, title, Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.65),
         size=30, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, subtitle, Inches(0.3), Inches(0.72), Inches(12.7), Inches(0.4),
             size=14, color=C_LIGHT, italic=True)


def _label(slide, text, l, t, w, h, color=C_ACCENT, size=11, bold=True):
    _txt(slide, text, l, t, w, h, size=size, bold=bold, color=color,
         align=PP_ALIGN.CENTER)


def _divider(slide, y):
    _rect(slide, Inches(0.3), y, Inches(12.7), Pt(1), bg=C_ACCENT)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def s_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _rect(slide, 0, 0, Inches(0.35), H, bg=C_ACCENT)
    _rect(slide, Inches(0.35), Inches(2.4), Inches(12.98), Inches(3.0), bg=C_PANEL)

    _txt(slide, "TESTTHEKEY", Inches(0.6), Inches(2.55), Inches(12.0), Inches(1.0),
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, "Architecture & Code Structure",
         Inches(0.6), Inches(3.45), Inches(12.0), Inches(0.6),
         size=22, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)
    _txt(slide, "LLM-Judge Fake Receipt Detector  |  Python  |  Streamlit  |  OpenAI  |  Langfuse",
         Inches(0.6), Inches(4.05), Inches(12.0), Inches(0.4),
         size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    _mono(slide,
          "testthekey/\n"
          "  streamlit_app.py       # 3-tab Streamlit UI\n"
          "  src/\n"
          "    config.py            # Settings + env loading\n"
          "    dataset.py           # Label loading, OCR parsing\n"
          "    features.py          # Blur, brightness, contrast\n"
          "    judges.py            # Prompt builder + run_judge()\n"
          "    vote.py              # majority_vote(), vote_tally()\n"
          "    langfuse_mcp_client.py  # HTTP/MCP client\n"
          "  scripts/               # CLI tools\n"
          "  tests/                 # 69 unit + integration tests",
          Inches(0.6), Inches(4.7), Inches(12.0), Inches(2.55),
          size=12, color=C_ACCENT)


# ── Slide 2 — End-to-end data flow ───────────────────────────────────────────

def s_dataflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "End-to-End Data Flow",
            "How a receipt image travels through the system")

    diagram = """\
Receipt image (.png/.jpg)
        |
        v
+-----------------------------------------------+
|           streamlit_app.py                    |
|   _browse_field()  ->  labels_df (CSV)        |
|   image_dir / receipt_id  ->  img_path        |
+-----------------------------------------------+
        |
        v
+---------------------------+    +---------------------------+
|  src/features.py          |    |  src/dataset.py           |
|                           |    |                           |
|  blur_variance_of_        |    |  image_basic_stats()      |
|    laplacian(img)         |    |    width, height,         |
|  brightness_contrast(img) |    |    file_kb, aspect_ratio  |
|    brightness_mean        |    |  extract_total_from_text()|
|    contrast_std           |    |    -> ocr_total (float)   |
+---------------------------+    +---------------------------+
        |                                 |
        +----------------+----------------+
                         |
                         v
        +--------------------------------+
        |  src/judges.py                 |
        |                                |
        |  JudgeConfig(name, model,      |
        |    temperature, persona)       |
        |                                |
        |  run_judge(client, cfg,        |
        |    image_path, receipt_id)     |
        |    -> (parsed, meta)           |
        |                                |
        |  parsed = {label, confidence,  |
        |    reasons[], flags[]}         |
        |  meta   = {usage, input,       |
        |    output}                     |
        +--------------------------------+
                         |   x3 judges
                         v
        +--------------------------------+
        |  src/vote.py                   |
        |                                |
        |  majority_vote(labels[])       |
        |    -> "FAKE"|"REAL"|"UNCERTAIN"|
        |  vote_tally(labels[])          |
        |    -> {FAKE:n, REAL:n, UNC:n}  |
        +--------------------------------+
                         |
                         v
        +--------------------------------+
        |  src/langfuse_mcp_client.py    |
        |  LangfuseMCPClient             |
        |  -> log_observation()          |
        |  -> log_generation() x3        |
        |  -> create_score()  x8         |
        |  -> dataset_run_log()          |
        +--------------------------------+"""

    _mono(slide, diagram, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.0),
          size=9.5, color=C_GREEN)


# ── Slide 3 — Module dependency map ──────────────────────────────────────────

def s_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Module Dependency Map",
            "What imports what — import graph of the project")

    diagram = """\
                    streamlit_app.py
                          |
          +---------------+------------------+------------------+
          |               |                  |                  |
          v               v                  v                  v
    src/config.py   src/dataset.py    src/features.py    src/judges.py
    load_settings() image_basic_       blur_variance_    JudgeConfig
    Settings        stats()            of_laplacian()    run_judge()
    (dataclass)     extract_total_     brightness_       _base_prompt()
                    from_text()        contrast()        _normalize_output()
                          |                                     |
                          |              +----------------------+
                          |              |
                          v              v
                    [PIL/Pillow]    [openai.OpenAI]
                    (image I/O)    (API calls)

          +---------------------------------------------------+
          |               src/vote.py                         |
          |   majority_vote(labels: list[str]) -> str         |
          |   vote_tally(labels: list[str]) -> dict           |
          |   (no external deps — pure Python)                |
          +---------------------------------------------------+

          +---------------------------------------------------+
          |          src/langfuse_mcp_client.py               |
          |   LangfuseMCPClient                               |
          |   imports: httpx, json  (stdlib)                  |
          |   talks to: localhost:8005/mcp/  via JSON-RPC     |
          |   no langfuse SDK, no credentials in app          |
          +---------------------------------------------------+

  External dependencies:
    openai  ->  GPT-4o / GPT-4.1-mini / GPT-4o-mini  (vision API)
    httpx   ->  MCP server HTTP transport
    plotly  ->  Charts in Dataset Statistics + Batch Evaluation tabs
    pandas  ->  DataFrame operations on label CSV
    PIL     ->  Image reading for feature extraction"""

    _mono(slide, diagram, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.0),
          size=10, color=C_ACCENT)


# ── Slide 4 — src/judges.py ───────────────────────────────────────────────────

def s_judges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "src/judges.py  — The Judge Engine",
            "Prompt construction, API call, output normalisation")

    left_code = """\
# Data model
@dataclass(frozen=True)
class JudgeConfig:
    name:        str    # "judge_1"
    model:       str    # "gpt-4o-mini"
    temperature: float  # 0.2
    persona:     str    # "strict, skeptical..."

# Entry point
def run_judge(
    client:     OpenAI,
    cfg:        JudgeConfig,
    image_path: Path,
    receipt_id: str
) -> tuple[dict, dict]:
    # 1. Build prompt with persona injected
    prompt   = _base_prompt(cfg.persona)
    # 2. Encode image as base64 data URL
    data_url = _image_to_data_url(image_path)
    # 3. Call OpenAI with vision content
    resp = client.chat.completions.create(
        model=cfg.model,
        temperature=cfg.temperature,
        messages=[system_msg, user_msg_with_image],
        response_format={"type": "json_object"},
    )
    # 4. Parse + normalise
    parsed = _normalize_output(
               _safe_parse_json(resp.choices[0].message.content))
    return parsed, {usage, input_payload, output_payload}"""

    right_code = """\
# Return shape
parsed = {
  "label":      "FAKE"|"REAL"|"UNCERTAIN",
  "confidence": float(0-100),
  "reasons":    list[str],  # 1-5 items
  "flags":      list[str],  # optional tags
}

meta = {
  "usage": {
    "input_tokens":  int,
    "output_tokens": int,
    "total_tokens":  int,
  },
  "input": {
    "receipt_id":   str,
    "prompt":       str,  # full prompt text
    "model":        str,
    "temperature":  float,
    "persona":      str,
    "image_path":   str,
  },
  "output": {
    "raw":    str,   # raw LLM response
    "parsed": dict,  # normalised
  },
}

# Fallback chain
# 1. json_object response_format
# 2. Extract first {...} block
# 3. Return {"label":"UNCERTAIN",...}"""

    _mono(slide, left_code,  Inches(0.3),  Inches(1.25), Inches(6.35), Inches(6.0), size=10)
    _mono(slide, right_code, Inches(6.75), Inches(1.25), Inches(6.25), Inches(6.0), size=10,
          color=C_YELLOW)


# ── Slide 5 — src/vote.py ────────────────────────────────────────────────────

def s_vote(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "src/vote.py  — Majority Vote Logic",
            "Pure Python aggregation — no external deps")

    code = """\
def majority_vote(labels: list[str]) -> str:
    \"\"\"
    Return the label that appears >= 2 times out of 3.
    Falls back to UNCERTAIN if no majority exists.
    \"\"\"
    counts = Counter(labels)
    for label in ("FAKE", "REAL"):
        if counts[label] >= 2:
            return label
    return "UNCERTAIN"


def vote_tally(labels: list[str]) -> dict[str, int]:
    \"\"\"Return raw vote counts for all three outcomes.\"\"\"
    base = {"FAKE": 0, "REAL": 0, "UNCERTAIN": 0}
    base.update(Counter(labels))
    return base"""

    examples = """\
Decision table — all 8 possible input combinations:

  labels input             majority_vote()   vote_tally()
  -------------------------------------------------------
  [FAKE, FAKE, FAKE]    -> FAKE              {F:3 R:0 U:0}
  [FAKE, FAKE, REAL]    -> FAKE              {F:2 R:1 U:0}
  [FAKE, FAKE, UNCERTAIN]-> FAKE             {F:2 R:0 U:1}
  [REAL, REAL, REAL]    -> REAL              {F:0 R:3 U:0}
  [REAL, REAL, FAKE]    -> REAL              {F:1 R:2 U:0}
  [REAL, REAL, UNCERTAIN]-> REAL             {F:0 R:2 U:1}
  [FAKE, REAL, UNCERTAIN]-> UNCERTAIN        {F:1 R:1 U:1}
  [UNC,  UNC,  UNC]     -> UNCERTAIN        {F:0 R:0 U:3}

  Priority: FAKE checked before REAL (conservative — flags first).
  UNCERTAIN is the fallback — never has >= 2 votes by itself
  unless all three judges return it."""

    _mono(slide, code,     Inches(0.3),  Inches(1.25), Inches(5.8), Inches(4.5), size=11)
    _mono(slide, examples, Inches(6.25), Inches(1.25), Inches(6.75), Inches(4.5),
          size=10.5, color=C_YELLOW)

    note = """\
Why majority vote over confidence-weighted average?
  Simple    — no tuning, no weights to calibrate
  Robust    — one outlier judge cannot swing the result
  Interpretable — you can see exactly which judges voted what
  UNCERTAIN is explicit signal, not a hidden 0.5 confidence"""
    _mono(slide, note, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.35),
          size=11, color=C_GREEN)


# ── Slide 6 — src/langfuse_mcp_client.py ─────────────────────────────────────

def s_mcp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "src/langfuse_mcp_client.py  — MCP HTTP Client",
            "Thin wrapper over the Langfuse MCP server — no SDK, no credentials")

    class_diag = """\
class LangfuseMCPClient:
    url:        str  = "http://localhost:8005/mcp/"
    timeout:    int  = 60
    _session_id: str | None

    # Transport
    _initialize()         # MCP handshake, captures session ID
    _call(method, params) # JSON-RPC POST, parses SSE or JSON response
    call_tool(name, args) # Calls tools/call, returns text string

    # Auth
    auth_check(verbose)   -> dict

    # Trace management
    create_trace_id(seed) -> str (deterministic UUID)
    get_trace(trace_id)   -> dict
    list_traces(name, limit)-> dict

    # Observation logging
    log_observation(name, as_type, trace, observation) -> dict
    log_generation(observation, trace, latency_ms)     -> dict
    log_batch(records)                                 -> dict
    log_error_event(message, trace, error_type, data)  -> dict

    # Scores
    create_score(trace_id, name, value, comment, metadata) -> dict

    # Dataset management
    dataset_create(name, description, metadata)       -> dict
    dataset_add_item(dataset_name, input,             -> dict
                     expected_output, item_id, metadata)
    dataset_run_log(run_name, dataset_item_id,        -> dict
                    trace_id, observation_id, metadata)

    # Prompts
    prompt_create_version(name, prompt, labels, tags) -> dict
    prompt_get(prompt_name, label)                    -> dict

    flush()               # langfuse_flush MCP tool"""

    protocol = """\
JSON-RPC / MCP Streamable HTTP protocol:

  POST /mcp/  (initialize)
  <-- mcp-session-id: <uuid>   (optional)

  POST /mcp/  (tools/call)
  --> {jsonrpc, id, method, params}
  <-- text/event-stream OR application/json

  Response parsing chain:
    1. application/json  -> resp.json()["result"]
    2. text/event-stream -> find "data: {...}" line
    3. Last resort       -> json.loads(resp.text)["result"]
    4. Raise RuntimeError with body preview

  Both POST calls use follow_redirects=True
  mcp-session-id header is optional (only sent if received)"""

    _mono(slide, class_diag, Inches(0.3),  Inches(1.25), Inches(7.3), Inches(6.0),
          size=9.5, color=C_ACCENT)
    _mono(slide, protocol,   Inches(7.65), Inches(1.25), Inches(5.35), Inches(6.0),
          size=10, color=C_YELLOW)


# ── Slide 7 — streamlit_app.py structure ─────────────────────────────────────

def s_streamlit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "streamlit_app.py  — Application Structure",
            "Module-level layout: constants, helpers, tab functions, main()")

    diagram = """\
streamlit_app.py
|
+-- IMPORTS
|     openai, streamlit, pandas, plotly, pathlib
|     src.config, src.dataset, src.features, src.judges
|     src.vote, src.langfuse_mcp_client, src.langfuse_logger
|
+-- MODULE-LEVEL CONSTANTS
|     _LABEL_COLORS  = {"REAL": "#2ecc71", "FAKE": "#e74c3c"}
|     _PLOTLY_THEME  = "plotly_white"
|     _MODEL_PRICING = {model: {input: $/1M, output: $/1M}, ...}
|     _TOKENS_PER_IMAGE_INPUT  = 1_500   # estimate
|     _TOKENS_PER_IMAGE_OUTPUT = 180     # estimate
|
+-- @st.cache_data
|     _compute_dataset_stats(label_path, image_dir, ocr_dir)
|       -> pd.DataFrame  [per-image features + label]
|
+-- HELPER FUNCTIONS
|     _estimate_cost(n_images, models)   -> list[dict]   # cost table rows
|     _hist(df, col, title, xlabel, by_label)  -> go.Figure
|     _pick_file(title, filetypes)       -> str   # tkinter file dialog
|     _pick_dir(title)                  -> str   # tkinter dir dialog
|     _browse_field(label, key, default, btn_key, is_dir) -> str
|
+-- TAB FUNCTIONS
|     _stats_tab(label_path, image_dir, ocr_dir)
|       Section 1 — REAL vs FAKE bar chart
|       Section 2 — receipt total histograms (if OCR)
|       Section 3 — image dimensions + file size
|       Section 4 — aspect ratio
|       Section 5 — blur / brightness / contrast
|       Section 6 — summary stats table by label
|
|     _batch_results_charts(results: list[dict])
|       -> summary metrics (5 st.metric)
|       -> confusion matrix (px.imshow)
|       -> prediction distribution (px.bar)
|       -> confidence box plots (px.box)
|       -> confidence histograms per judge
|       -> per-judge accuracy table
|       -> wrong predictions table
|
|     _batch_tab(labels_df, image_dir, ocr_dir, client, judge_cfgs, settings)
|       -> pool filter radio + langfuse checkbox
|       -> quick-select: slider + Random/All/Clear buttons
|       -> run name text_input (if langfuse)
|       -> st.multiselect (key="batch_multiselect")
|       -> thumbnail grid (6 cols, max 30)
|       -> cost estimation table
|       -> Run / Clear buttons
|       -> evaluation loop  [run_judge x3 + vote + langfuse]
|       -> _batch_results_charts()  +  JSON download button
|
+-- main()
      load_dotenv(), load_settings(), OpenAI(api_key)
      sidebar: dataset mode selector + browse fields
      sidebar: label CSV read -> labels_df
      sidebar: receipt picker (selectbox)
      sidebar: Langfuse dataset registration section
      judge_cfgs = [JudgeConfig x3]  <- shared across tabs
      tab1, tab2, tab3 = st.tabs([...])
        tab1 -> single receipt analysis (+ langfuse_logger)
        tab2 -> _stats_tab()
        tab3 -> _batch_tab()"""

    _mono(slide, diagram, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.0),
          size=9, color=C_ACCENT)


# ── Slide 8 — Session state & Streamlit re-run flow ──────────────────────────

def s_session(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Streamlit Session State & Re-run Flow",
            "How state is preserved across interactions")

    diagram = """\
st.session_state keys used by the app:
+---------------------------+------------------+------------------------------+
| Key                       | Type             | Set by                       |
+---------------------------+------------------+------------------------------+
| _label_path__val          | str (path)       | _browse_field() + Browse btn |
| _image_dir__val           | str (path)       | _browse_field() + Browse btn |
| _ocr_dir__val             | str (path)       | _browse_field() + Browse btn |
| batch_multiselect         | list[str]        | Quick-select buttons + user  |
| batch_lf_run_name         | str              | st.text_input (key=)         |
| batch_results             | list[dict]       | Evaluation loop on run_click |
+---------------------------+------------------+------------------------------+

Browse button pattern (avoids key conflict with st.text_input):
  backing_key = f"{state_key}__val"      # e.g. "_label_path__val"
  st.session_state[backing_key] = typed  # keep in sync when user types
  if browse_clicked:
      st.session_state[backing_key] = picked_path
      st.rerun()                         # triggers immediate re-render

Quick-select button pattern (sets widget state before render):
  if random_clicked:
      st.session_state["batch_multiselect"] = sampled_list
      st.rerun()                         # multiselect reads from session state

Re-run trigger chain:
  User edits field  -> Streamlit auto re-runs   (no rerun() needed)
  Browse button     -> _browse_field() calls st.rerun()
  Quick-select btn  -> _batch_tab() calls st.rerun()
  Run batch eval    -> loop finishes -> st.session_state["batch_results"] = ...
                    -> st.rerun()    -> results section renders
  Clear results     -> del st.session_state["batch_results"]
                    -> st.rerun()

@st.cache_data usage:
  _compute_dataset_stats(label_path, image_dir, ocr_dir)
    Keyed by all three string args.
    Re-computed only when dataset path changes.
    Scanning 218 images takes ~5-15 seconds — cached after first run."""

    _mono(slide, diagram, Inches(0.3), Inches(1.25), Inches(12.7), Inches(5.95),
          size=10, color=C_YELLOW)


# ── Slide 9 — Langfuse logging trace structure ────────────────────────────────

def s_langfuse_trace(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Langfuse Trace Structure  — Per Receipt",
            "What one batch evaluation trace looks like in Langfuse")

    diagram = """\
TRACE  batch_eval_<receipt_id>
  input:  {image_id, ground_truth, width, height, file_kb,
           brightness_mean, contrast_std, blur_variance}
  output: {final_label, tally, is_correct, judge_labels[]}
  tags:   [receipt, batch-eval, real|fake, <run_name>]
  |
  +-- SPAN  dataset_analysis
  |     input:  {image_id, ground_truth}
  |     output: {width, height, file_kb, aspect_ratio,
  |               brightness_mean, contrast_std, blur_variance}
  |
  +-- GENERATION  judge_1   (gpt-4o-mini, temp=0.2)
  |     input:  {system: <full prompt with persona>,
  |               user: "Receipt ID: X123.png [image attached]"}
  |     output: {label, confidence, reasons[], flags[]}
  |     usage:  {input: N tokens, output: M tokens, unit: TOKENS}
  |     metadata: {temperature, persona, ground_truth,
  |                latency_ms, raw_response[:300]}
  |
  +-- GENERATION  judge_2   (gpt-4.1-mini, temp=0.4)
  |     [same structure]
  |
  +-- GENERATION  judge_3   (gpt-4o, temp=0.7)
  |     [same structure]
  |
  +-- SPAN  vote_aggregation
        input:  {judge_labels: [FAKE, REAL, REAL]}
        output: {final_label: REAL, tally: {F:1,R:2,U:0},
                  is_correct: true}

SCORES  (8 per trace, all attached to trace_id):
  judge_1_confidence   float  e.g. 62.0   comment: label=UNCERTAIN
  judge_1_correctness  float  1.0/0.5/0.0 comment: judge=UNC gt=REAL
  judge_2_confidence   float  e.g. 85.0
  judge_2_correctness  float  1.0/0.5/0.0
  judge_3_confidence   float  e.g. 91.0
  judge_3_correctness  float  1.0/0.5/0.0
  final_correct        float  1.0 or 0.0   comment: final=REAL gt=REAL
  inter_judge_agreement float  1.0 or 0.0   comment: labels=[F,R,R]

DATASET RUN LOG:
  run_name:        "batch-eval-v1"
  dataset_item_id: "X51005230616.png"
  trace_id:        <actual Langfuse UUID from log_observation response>"""

    _mono(slide, diagram, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.0),
          size=9.5, color=C_GREEN)


# ── Slide 10 — src/config.py + settings flow ─────────────────────────────────

def s_config(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Configuration & Settings Flow",
            "src/config.py  +  .env  ->  Settings dataclass  ->  all modules")

    left = """\
# .env  (never committed)
OPENAI_API_KEY=sk-...
JUDGE_MODELS=gpt-4o-mini,gpt-4.1-mini,gpt-4o
LANGFUSE_PUBLIC_KEY=pk-...  (optional)
LANGFUSE_SECRET_KEY=sk-...  (optional)
LANGFUSE_HOST=http://...    (optional)

# src/config.py
@dataclass
class Settings:
    openai_api_key:    str
    judge_models:      list[str]  # [model1, model2, model3]
    langfuse_public_key: str | None
    langfuse_secret_key: str | None
    langfuse_host:      str | None

def load_settings() -> Settings:
    return Settings(
        openai_api_key  = os.getenv("OPENAI_API_KEY"),
        judge_models    = os.getenv(
            "JUDGE_MODELS",
            "gpt-4o-mini,gpt-4.1-mini,gpt-4o"
        ).split(","),
        langfuse_public_key = os.getenv("LANGFUSE_PUBLIC_KEY"),
        langfuse_secret_key = os.getenv("LANGFUSE_SECRET_KEY"),
        langfuse_host       = os.getenv("LANGFUSE_HOST"),
    )"""

    right = """\
# Usage in main()
settings = load_settings()
client   = OpenAI(api_key=settings.openai_api_key)

judge_cfgs = [
    JudgeConfig(
        name="judge_1",
        model=settings.judge_models[0],  # gpt-4o-mini
        temperature=0.2,
        persona="strict, skeptical...",
    ),
    JudgeConfig(
        name="judge_2",
        model=settings.judge_models[1],  # gpt-4.1-mini
        temperature=0.4,
        persona="balanced...",
    ),
    JudgeConfig(
        name="judge_3",
        model=settings.judge_models[2],  # gpt-4o
        temperature=0.7,
        persona="lenient...",
    ),
]
# judge_cfgs is defined BEFORE tabs
# so Tab 1 and Tab 3 share the same list

# LangfuseMCPClient does NOT use settings —
# credentials live inside the MCP server
lf = LangfuseMCPClient()  # no keys needed"""

    _mono(slide, left,  Inches(0.3),  Inches(1.25), Inches(6.3), Inches(6.0), size=10)
    _mono(slide, right, Inches(6.7),  Inches(1.25), Inches(6.3), Inches(6.0),
          size=10, color=C_YELLOW)


# ── Slide 11 — Batch evaluation loop code ────────────────────────────────────

def s_batch_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Batch Evaluation Loop  — Core Logic",
            "_batch_tab() inner loop: judges + vote + Langfuse per image")

    code = """\
for idx, (_, row) in enumerate(selected_df.iterrows()):
    img_id   = str(row["image"])
    gt       = str(row["label"]).upper().strip()
    img_path = image_dir / img_id

    prog.progress((idx+1)/n_batch, text=f"{idx+1}/{n_batch} {img_id}")

    labels:          list[str]  = []
    judge_outputs:   list[dict] = []
    judge_metas:     list[dict] = []
    judge_latencies: list[int]  = []

    try:
        for cfg in judge_cfgs:                    # 3 judges
            t0 = time.time()
            parsed, jmeta = run_judge(            # -> API call
                client, cfg,
                image_path=img_path,
                receipt_id=img_id,
            )
            judge_latencies.append(               # real ms
                int((time.time() - t0) * 1000))
            labels.append(parsed["label"])
            judge_outputs.append({
                "name": cfg.name, "model": cfg.model, **parsed})
            judge_metas.append(jmeta)

        final_label = majority_vote(labels)       # FAKE/REAL/UNC
        tally       = vote_tally(labels)          # {F:n, R:n, U:n}
        is_correct  = final_label == gt

        results.append({                          # local cache
            "image_id": img_id, "ground_truth": gt,
            "final_label": final_label,
            "tally": tally, "is_correct": is_correct,
            "judges": judge_outputs,
        })

        if lf is not None:
            # ... comprehensive Langfuse logging block ...
            # (analysis span + 3 generations + vote span
            #  + 8 scores + dataset_run_log)
            # All errors caught silently — batch never aborted

    except Exception as exc:
        results.append({                          # error row
            "image_id": img_id, "error": str(exc),
            "final_label": "ERROR", ...
        })

st.session_state["batch_results"] = results
st.rerun()  # triggers results section to render"""

    _mono(slide, code, Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.05),
          size=10, color=C_GREEN)


# ── Slide 12 — Testing strategy ───────────────────────────────────────────────

def s_tests(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header(slide, "Testing Strategy",
            "tests/  —  69 tests, zero real API calls for unit suite")

    layout = """\
tests/
  test_vote.py             # 15 unit tests — pure logic
  test_dataset.py          # 12 unit tests — CSV loading, OCR parsing
  test_features.py         # 10 unit tests — blur/brightness/contrast
  test_judges.py           # 18 unit tests — prompt builder, normalisation
  test_langfuse_client.py  #  6 unit tests — MCP response parsing
  test_integration.py      #  8 integration tests (real API, marked)

Run unit tests (no API cost, < 5 seconds):
  pytest tests/ -v -m "not integration"

Run all tests including API calls (~$0.01):
  pytest tests/ -v"""

    vote_tests = """\
# test_vote.py examples
def test_fake_majority():
    assert majority_vote(["FAKE","FAKE","REAL"]) == "FAKE"

def test_uncertain_on_split():
    assert majority_vote(["FAKE","REAL","UNCERTAIN"]) == "UNCERTAIN"

def test_tally_keys():
    t = vote_tally(["FAKE","REAL","REAL"])
    assert t == {"FAKE":1,"REAL":2,"UNCERTAIN":0}"""

    judge_tests = """\
# test_judges.py — mock OpenAI
@pytest.fixture
def mock_client(monkeypatch):
    resp = MagicMock()
    resp.choices[0].message.content = json.dumps({
        "label":"FAKE","confidence":85,
        "reasons":["font mismatch"],"flags":[]})
    resp.usage.prompt_tokens    = 1200
    resp.usage.completion_tokens = 80
    resp.usage.total_tokens     = 1280
    client = MagicMock()
    client.chat.completions.create.return_value = resp
    return client

def test_run_judge_label(mock_client, tmp_image):
    parsed, meta = run_judge(
        mock_client, cfg, tmp_image, "test.png")
    assert parsed["label"] == "FAKE"
    assert 0 <= parsed["confidence"] <= 100
    assert isinstance(parsed["reasons"], list)"""

    _mono(slide, layout,      Inches(0.3),  Inches(1.25), Inches(4.5),  Inches(3.1), size=10.5)
    _mono(slide, vote_tests,  Inches(0.3),  Inches(4.5),  Inches(4.5),  Inches(2.75), size=10)
    _mono(slide, judge_tests, Inches(4.95), Inches(1.25), Inches(8.05), Inches(5.95), size=9.5,
          color=C_YELLOW)


# ── Slide 13 — End / summary ──────────────────────────────────────────────────

def s_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _rect(slide, 0, 0, Inches(0.35), H, bg=C_ACCENT)
    _rect(slide, Inches(0.35), Inches(1.8), Inches(12.98), Inches(4.2), bg=C_PANEL)

    _txt(slide, "Architecture Summary",
         Inches(0.6), Inches(2.0), Inches(12.0), Inches(0.75),
         size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    summary = """\
Receipt image
      |
      v
src/features.py + src/dataset.py     CPU-only feature extraction (free)
      |
      v
src/judges.py  run_judge() x3        OpenAI Vision API — parallel LLM calls
      |
      v
src/vote.py  majority_vote()         Pure Python — no external deps
      |
      +---> Local results cache      JSON in session_state / eval_cache.json
      |
      +---> src/langfuse_mcp_client  HTTP/JSON-RPC to MCP server
              LangfuseMCPClient      5 logging calls per receipt image
                                     Full observability — zero app credentials"""

    _mono(slide, summary, Inches(1.2), Inches(2.85), Inches(11.0), Inches(2.9),
          size=12, color=C_ACCENT, bg=C_CODE)

    _txt(slide, "TESTTHEKEY  |  Architecture & Code Structure Presentation",
         Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.4),
         size=13, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = _prs()

    s_title(prs)
    s_dataflow(prs)
    s_modules(prs)
    s_judges(prs)
    s_vote(prs)
    s_mcp(prs)
    s_streamlit(prs)
    s_session(prs)
    s_langfuse_trace(prs)
    s_config(prs)
    s_batch_loop(prs)
    s_tests(prs)
    s_end(prs)

    out = Path("TESTTHEKEY_Architecture.pptx")
    prs.save(str(out))
    print(f"Saved: {out.resolve()}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
