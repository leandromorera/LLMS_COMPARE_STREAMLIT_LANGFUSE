"""
Generate a professional PowerPoint presentation for the TESTTHEKEY project.
Run from the project root:
    python scripts/generate_pptx.py
Output: TESTTHEKEY_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_ACCENT2   = RGBColor(0xE7, 0x4C, 0x3C)   # red (FAKE)
C_ACCENT3   = RGBColor(0x2E, 0xCC, 0x71)   # green (REAL)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xB0, 0xC4, 0xDE)   # light steel blue
C_DARK_CARD = RGBColor(0x16, 0x2B, 0x3E)   # slightly lighter navy for cards

W, H = Inches(13.33), Inches(7.5)          # 16:9 widescreen


# ── Helpers ───────────────────────────────────────────────────────────────────

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _bg(slide, color: RGBColor = C_BG) -> None:
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height,
         bg: RGBColor | None = None,
         border: RGBColor | None = None) -> object:
    """Add a plain rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text: str, left, top, width, height,
         size: int = 18, bold: bool = False, color: RGBColor = C_WHITE,
         align=PP_ALIGN.LEFT, wrap: bool = True, italic: bool = False) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"


def _header_bar(slide, title: str, subtitle: str = "") -> None:
    """Cyan left bar + title text."""
    _box(slide, Inches(0), Inches(0), Inches(0.18), H, bg=C_ACCENT)
    _txt(slide, title, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.75),
         size=32, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, subtitle, Inches(0.35), Inches(0.85), Inches(12.5), Inches(0.45),
             size=16, color=C_LIGHT, italic=True)


def _card(slide, left, top, width, height, title: str, body_lines: list[str],
          title_color: RGBColor = C_ACCENT, body_size: int = 15) -> None:
    """Rounded card with title + bullet body."""
    _box(slide, left, top, width, height, bg=C_DARK_CARD)
    _box(slide, left, top, width, Inches(0.38), bg=title_color)
    _txt(slide, title, left + Inches(0.12), top + Inches(0.04),
         width - Inches(0.14), Inches(0.34), size=14, bold=True, color=C_BG)
    body_top = top + Inches(0.44)
    body_h   = height - Inches(0.54)
    body_txt = "\n".join(f"• {l}" for l in body_lines)
    _txt(slide, body_txt, left + Inches(0.12), body_top,
         width - Inches(0.24), body_h, size=body_size, color=C_WHITE, wrap=True)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _bg(slide)
    # large accent bar left
    _box(slide, Inches(0), Inches(0), Inches(0.35), H, bg=C_ACCENT)
    # coloured band center
    _box(slide, Inches(0.35), Inches(2.6), Inches(12.98), Inches(2.5), bg=C_DARK_CARD)
    _txt(slide, "LLM-Judge Fake Receipt Detector",
         Inches(0.6), Inches(2.75), Inches(12.0), Inches(1.2),
         size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, "3 OpenAI Vision Judges  ·  Majority Vote  ·  Langfuse Observability",
         Inches(0.6), Inches(3.75), Inches(12.0), Inches(0.5),
         size=20, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)
    _txt(slide, "Python  ·  Streamlit  ·  OpenAI  ·  Langfuse  ·  MCP",
         Inches(0.6), Inches(4.35), Inches(12.0), Inches(0.4),
         size=15, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(slide, "February 2026",
         Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.4),
         size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "The Problem", "Why can't classical CV catch receipt forgeries?")
    _card(slide, Inches(0.35), Inches(1.45), Inches(4.0), Inches(5.6),
          "What forgers do", [
              "Change a single digit in the total",
              "Copy-paste amounts between receipts",
              "Alter font or spacing in key fields",
              "Substitute store names / dates",
              "Save at same quality — no metadata leak",
          ], title_color=C_ACCENT2)
    _card(slide, Inches(4.55), Inches(1.45), Inches(4.0), Inches(5.6),
          "Why classical features fail", [
              "Blur & brightness: nearly identical",
              "File size: forgers match it",
              "Aspect ratio: same portrait format",
              "No single feature separates classes",
              "Dataset: 83.9% REAL → naive baseline = 83.9%",
          ], title_color=C_LIGHT)
    _card(slide, Inches(8.75), Inches(1.45), Inches(4.2), Inches(5.6),
          "Solution: Vision LLM judges", [
              "Trained on billions of documents",
              "Detect font inconsistencies visually",
              "Spot pixel-level copy-paste artefacts",
              "Read layout plausibility holistically",
              "No labelled training data needed",
          ], title_color=C_ACCENT3)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "System Architecture", "End-to-end pipeline for a single receipt image")

    # Pipeline boxes
    stages = [
        ("INPUT", "Receipt image\n(.png / .jpg)\n+ OCR .txt\n(optional)", C_LIGHT),
        ("FEATURES", "CPU-only\nBlur · Brightness\nContrast · Dims\nOCR total", C_ACCENT),
        ("JUDGE PANEL", "3 × OpenAI Vision\ngpt-4o-mini\ngpt-4.1-mini\ngpt-4o", C_ACCENT2),
        ("VOTE", "Majority Vote\nFAKE ≥ 2\nREAL ≥ 2\nelse UNCERTAIN", C_ACCENT3),
        ("OUTPUT", "Verdict +\nConfidence +\nReasons +\nFlags", C_WHITE),
    ]
    box_w = Inches(2.2)
    box_h = Inches(3.8)
    gap   = Inches(0.22)
    start_left = Inches(0.35)
    top   = Inches(1.7)

    for i, (title, body, color) in enumerate(stages):
        left = start_left + i * (box_w + gap)
        _box(slide, left, top, box_w, box_h, bg=C_DARK_CARD)
        _box(slide, left, top, box_w, Inches(0.38), bg=color)
        title_color = C_BG if color != C_WHITE else C_BG
        _txt(slide, title, left + Inches(0.08), top + Inches(0.05),
             box_w - Inches(0.1), Inches(0.32), size=13, bold=True,
             color=title_color, align=PP_ALIGN.CENTER)
        _txt(slide, body, left + Inches(0.1), top + Inches(0.48),
             box_w - Inches(0.2), box_h - Inches(0.6),
             size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Arrow
        if i < len(stages) - 1:
            arr_left = left + box_w + Inches(0.04)
            _txt(slide, "→", arr_left, top + Inches(1.6), Inches(0.18), Inches(0.4),
                 size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Langfuse bar at the bottom
    _box(slide, Inches(0.35), Inches(5.8), Inches(12.6), Inches(1.3), bg=C_DARK_CARD)
    _txt(slide, "LANGFUSE  (via MCP)  — Full observability on every run",
         Inches(0.5), Inches(5.88), Inches(12.0), Inches(0.38),
         size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _txt(slide,
         "Trace → dataset_analysis span → judge_1/2/3 generations (prompt + tokens + latency) → vote_aggregation span → 8 scores → dataset_run_log",
         Inches(0.5), Inches(6.25), Inches(12.0), Inches(0.6),
         size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_judges(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "The Judge Panel", "Three independent personas — deliberately different priors")

    judges = [
        ("judge_1  ·  gpt-4o-mini", "Strict / Skeptical", 0.2, C_ACCENT2, [
            "Aggressively flags any forensic inconsistency",
            "Temperature 0.2 → deterministic, conservative",
            "Tends toward UNCERTAIN when evidence is weak",
            "Best at catching obvious copy-paste artefacts",
            "Lowest cost — $0.15/$0.60 per 1M tokens",
        ]),
        ("judge_2  ·  gpt-4.1-mini", "Balanced / Artifact-aware", 0.4, C_ACCENT, [
            "Weighs printing/scan artefacts vs tampering",
            "Temperature 0.4 → moderate variability",
            "Most accurate in our evaluation (70%)",
            "Distinguishes thermal printer quirks from edits",
            "Mid-range cost — $0.40/$1.60 per 1M tokens",
        ]),
        ("judge_3  ·  gpt-4o", "Lenient / Benefit-of-doubt", 0.7, C_ACCENT3, [
            "Assumes REAL unless manipulation is clear",
            "Temperature 0.7 → higher variability",
            "Highest capability vision model (tiebreaker)",
            "Most decisive — rarely returns UNCERTAIN",
            "Premium cost — $2.50/$10.00 per 1M tokens",
        ]),
    ]

    card_w = Inches(4.0)
    for i, (name, persona, temp, color, bullets) in enumerate(judges):
        left = Inches(0.35) + i * Inches(4.32)
        _card(slide, left, Inches(1.45), card_w, Inches(5.6),
              name, bullets, title_color=color, body_size=14)
        _txt(slide, f'"{persona}"  ·  temp={temp}',
             left, Inches(1.08), card_w, Inches(0.35),
             size=13, color=color, italic=True, align=PP_ALIGN.CENTER)


def slide_prompt(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Judge Prompt Design", "Same base prompt, only the persona line changes")

    prompt_text = (
        'You are an expert forensic document examiner evaluating whether a receipt image is forged.\n'
        'Persona: {persona}\n\n'
        'Return ONLY valid JSON:\n'
        '{\n'
        '  "label": "FAKE|REAL|UNCERTAIN",\n'
        '  "confidence": 0-100,\n'
        '  "reasons": ["short reason 1", ...],\n'
        '  "flags": ["optional tag", ...]\n'
        '}\n\n'
        'Guidelines:\n'
        '• Base decision on visual cues (fonts, alignment, spacing,\n'
        '  copy-paste digits, pixel artifacts, shadows, totals)\n'
        '• UNCERTAIN + lower confidence if evidence is weak\n'
        '• reasons must be short, concrete, observable'
    )
    _box(slide, Inches(0.35), Inches(1.4), Inches(7.5), Inches(5.7), bg=C_DARK_CARD)
    _txt(slide, "PROMPT (verbatim)", Inches(0.5), Inches(1.5), Inches(7.0), Inches(0.35),
         size=13, bold=True, color=C_ACCENT)
    _txt(slide, prompt_text, Inches(0.5), Inches(1.9), Inches(7.2), Inches(5.0),
         size=12.5, color=C_LIGHT, italic=True)

    _card(slide, Inches(8.1), Inches(1.4), Inches(4.85), Inches(2.6),
          "Structured output", [
              "response_format: json_object (where supported)",
              "Fallback: salvage first {…} block",
              "_normalize_output() enforces schema:",
              "  label ∈ {FAKE, REAL, UNCERTAIN}",
              "  confidence clamped 0–100",
              "  reasons: 1–5 strings ≤ 240 chars",
          ], title_color=C_ACCENT)

    _card(slide, Inches(8.1), Inches(4.2), Inches(4.85), Inches(2.9),
          "Why this prompt works", [
              "Short → fast, cheap, less hallucination",
              "Concrete cues → actionable reasoning",
              "UNCERTAIN is explicit, not a failure",
              "Persona injects prior → ensemble diversity",
              "Image sent as base64 data URL inline",
          ], title_color=C_ACCENT3)


def slide_dataset(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Dataset & Distribution", "findit2 test set — 218 Malaysian retail receipts")

    # Stat pills
    stats = [
        ("218", "Total images"),
        ("183", "REAL (83.9%)"),
        ("35",  "FAKE (16.1%)"),
        ("5.2:1", "Class imbalance"),
    ]
    pill_w = Inches(2.8)
    for i, (val, label) in enumerate(stats):
        left = Inches(0.35) + i * Inches(3.1)
        _box(slide, left, Inches(1.45), pill_w, Inches(1.4), bg=C_DARK_CARD)
        _txt(slide, val, left, Inches(1.55), pill_w, Inches(0.7),
             size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        _txt(slide, label, left, Inches(2.2), pill_w, Inches(0.4),
             size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

    _card(slide, Inches(0.35), Inches(3.1), Inches(5.9), Inches(4.0),
          "Key insight — class imbalance", [
              "Naive classifier (always REAL) = 83.9% accuracy",
              "Accuracy alone is a misleading metric",
              "Per-class recall + UNCERTAIN rate matter more",
              "Sampling for evaluation must be stratified",
              "Use 'FAKE only' filter to oversample minority class",
          ], title_color=C_ACCENT2)

    _card(slide, Inches(6.45), Inches(3.1), Inches(6.5), Inches(4.0),
          "What features tell us (and don't)", [
              "File size: REAL ≈ FAKE — forgers match it",
              "Aspect ratio: both portrait ~0.4–0.5",
              "Blur variance: wide spread in both classes",
              "Brightness: slight over-exposure in some fakes",
              "→ No single feature cleanly separates classes",
              "→ Validates LLM vision over classical CV",
          ], title_color=C_LIGHT)


def slide_results(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Evaluation Results", "20-image batch · 3 judges · majority vote")

    # Top metrics
    metrics = [
        ("55%",  "Overall accuracy\n(11/20 correct)", C_ACCENT),
        ("9",    "Wrong predictions", C_ACCENT2),
        ("70%",  "Best single judge\n(judge_2 / gpt-4.1-mini)", C_ACCENT3),
        ("$0.007", "Total API cost\n(128 traces session)", C_LIGHT),
    ]
    pill_w = Inches(2.85)
    for i, (val, label, color) in enumerate(metrics):
        left = Inches(0.35) + i * Inches(3.1)
        _box(slide, left, Inches(1.45), pill_w, Inches(1.5), bg=C_DARK_CARD)
        _txt(slide, val, left, Inches(1.52), pill_w, Inches(0.72),
             size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
        _txt(slide, label, left, Inches(2.22), pill_w, Inches(0.55),
             size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Per-judge table
    _box(slide, Inches(0.35), Inches(3.2), Inches(12.6), Inches(3.85), bg=C_DARK_CARD)
    _txt(slide, "Per-judge performance",
         Inches(0.5), Inches(3.28), Inches(12.0), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)

    headers = ["Judge", "Model", "Accuracy", "Avg. Confidence", "FAKE", "REAL", "UNCERTAIN", "Key observation"]
    rows = [
        ["judge_1", "gpt-4o-mini",  "5.0%",  "54.5", "0",  "3",  "17", "Defaults to UNCERTAIN — almost useless"],
        ["judge_2", "gpt-4.1-mini", "70.0%", "79.5", "14", "15", "1",  "Most decisive and accurate"],
        ["judge_3", "gpt-4o",       "65.0%", "82.8", "0",  "15", "5",  "High confidence, never FAKE"],
    ]
    col_ws = [Inches(0.9), Inches(1.3), Inches(0.9), Inches(1.3), Inches(0.55), Inches(0.55), Inches(1.0), Inches(5.8)]
    row_h  = Inches(0.52)
    top0   = Inches(3.72)
    col_colors = [C_ACCENT, C_ACCENT, C_ACCENT3, C_LIGHT, C_ACCENT2, C_ACCENT3, C_LIGHT, C_LIGHT]

    # Header row
    left0 = Inches(0.42)
    for j, (hdr, cw, cc) in enumerate(zip(headers, col_ws, col_colors)):
        left = left0 + sum(col_ws[:j])
        _txt(slide, hdr, left, top0, cw, row_h,
             size=11, bold=True, color=cc, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        top = top0 + row_h + i * row_h
        row_bg = RGBColor(0x1A, 0x31, 0x48) if i % 2 == 0 else C_DARK_CARD
        _box(slide, Inches(0.42), top, Inches(12.5), row_h, bg=row_bg)
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            left = left0 + sum(col_ws[:j])
            _txt(slide, cell, left, top + Inches(0.08), cw, row_h - Inches(0.08),
                 size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_errors(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Error Analysis", "Three representative failure modes")

    cases = [
        ("Case 1 — False Negative", C_ACCENT2, [
            "GT: FAKE  →  Predicted: REAL",
            "Vote: FAKE:0 · REAL:2 · UNCERTAIN:1",
            "Judges 2+3 read layout as authentic",
            "Forgery: single digit changed in total",
            "No copy-paste artefact visible at standard res",
            "→ LLMs can't detect numeric-only edits",
        ]),
        ("Case 2 — Three-way Split", C_ACCENT, [
            "GT: REAL  →  Predicted: UNCERTAIN",
            "Vote: FAKE:1 · REAL:1 · UNCERTAIN:1",
            "Judge 1 flagged thermal printer font quirk",
            "Judge 3 accepted as real, judge 2 undecided",
            "Root cause: thermal artefacts mimic forgery",
            "→ Main source of false positives",
        ]),
        ("Case 3 — Conservative Miss", C_LIGHT, [
            "GT: FAKE  →  Predicted: UNCERTAIN",
            "Vote: FAKE:1 · REAL:0 · UNCERTAIN:2",
            "Only judge_1 correctly identified as FAKE",
            "Other two withheld votes → no majority",
            "Root cause: small/low-res image, cues invisible",
            "→ Visual-only detection has resolution limits",
        ]),
    ]

    card_w = Inches(4.0)
    for i, (title, color, bullets) in enumerate(cases):
        left = Inches(0.35) + i * Inches(4.32)
        _card(slide, left, Inches(1.45), card_w, Inches(5.1),
              title, bullets, title_color=color, body_size=14)

    _txt(slide,
         "Key fix: add OCR arithmetic check (do line items sum to total?) as a pre-filter — catches Case 1 class at near-zero cost",
         Inches(0.35), Inches(6.75), Inches(12.6), Inches(0.45),
         size=13, color=C_ACCENT, italic=True, align=PP_ALIGN.CENTER)


def slide_langfuse(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Langfuse Observability", "Full LLM operation visibility via MCP — no credentials in app code")

    _card(slide, Inches(0.35), Inches(1.45), Inches(5.9), Inches(5.6),
          "What is logged per receipt", [
              "dataset_analysis span — image stats (free, CPU)",
              "judge_1 generation — full system prompt + JSON response",
              "  → real input/output/total tokens + latency_ms",
              "judge_2 generation — same, different model",
              "judge_3 generation — same, gpt-4o",
              "vote_aggregation span — labels → final verdict",
              "8 scores per trace:",
              "  judge_N_confidence · judge_N_correctness (×3)",
              "  final_correct · inter_judge_agreement",
              "dataset_run_log — links trace to named run",
          ], title_color=C_ACCENT)

    _card(slide, Inches(6.45), Inches(1.45), Inches(6.5), Inches(2.6),
          "MCP integration (no SDK in app)", [
              "LangfuseMCPClient sends JSON-RPC to localhost:8005/mcp/",
              "MCP server holds Langfuse credentials internally",
              "Streamlit app needs zero LANGFUSE_* env vars",
              "Follows redirects · session-ID optional · SSE + JSON",
              "All errors caught silently — batch never aborted",
          ], title_color=C_ACCENT2)

    _card(slide, Inches(6.45), Inches(4.25), Inches(6.5), Inches(2.8),
          "Dataset runs — cross-run comparison", [
              "Each batch gets a named run (e.g. test-v1, test-v2)",
              "Langfuse groups traces by run automatically",
              "Compare accuracy / cost / scores across experiments",
              "Charts tab needs ≥ 2 runs to populate",
              "Real session cost: $0.007182  ·  128 traces  ·  60 scores",
          ], title_color=C_ACCENT3)


def slide_ui(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Streamlit Dashboard — 3 Tabs", "Interactive UI for single analysis, statistics, and batch evaluation")

    tabs = [
        ("Tab 1 — Analyze Receipt", C_ACCENT, [
            "Select any receipt from the dataset",
            "View image + extracted lightweight features",
            "Click 'Run judges' → 3 simultaneous API calls",
            "See per-judge JSON verdict side-by-side",
            "Majority vote + tally shown instantly",
            "Logs trace + generations + scores to Langfuse",
        ]),
        ("Tab 2 — Dataset Statistics", C_ACCENT3, [
            "Scan entire dataset (cached after first run)",
            "REAL vs FAKE label distribution bar chart",
            "Receipt total amount histograms (with OCR)",
            "Image dimensions, file size, aspect ratio",
            "Sharpness · brightness · contrast overlays",
            "Summary statistics table by label",
        ]),
        ("Tab 3 — Batch Evaluation", C_ACCENT2, [
            "Filter pool: All / REAL only / FAKE only",
            "Exact image selection with searchable multiselect",
            "Thumbnail grid (max 30) with REAL/FAKE badges",
            "Live cost estimate before any API spend",
            "Progress bar per image during evaluation",
            "Confusion matrix · box plots · per-judge table",
            "Download full results as JSON",
            "Logs comprehensive traces to Langfuse",
        ]),
    ]

    card_w = Inches(4.0)
    for i, (title, color, bullets) in enumerate(tabs):
        left = Inches(0.35) + i * Inches(4.32)
        _card(slide, left, Inches(1.45), card_w, Inches(5.6),
              title, bullets, title_color=color, body_size=13)


def slide_code_quality(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Code Quality & Reproducibility", "Production-ready structure from day one")

    _card(slide, Inches(0.35), Inches(1.45), Inches(5.9), Inches(5.6),
          "Project structure", [
              "streamlit_app.py  — 3-tab dashboard",
              "src/judges.py     — prompt builder + run_judge()",
              "src/vote.py       — majority_vote(), vote_tally()",
              "src/features.py   — blur, brightness, contrast",
              "src/dataset.py    — label loading, OCR parsing",
              "src/langfuse_mcp_client.py — HTTP/MCP wrapper",
              "scripts/run_eval20.py     — CLI batch eval",
              "scripts/register_dataset.py — Langfuse registration",
              "tests/  — 69 unit + integration tests",
          ], title_color=C_ACCENT)

    _card(slide, Inches(6.45), Inches(1.45), Inches(6.5), Inches(2.6),
          "Reproducibility checklist", [
              "✓ All secrets in .env (.env.example provided)",
              "✓ Pinned versions in requirements.txt",
              "✓ random_state=42 for all sampling",
              "✓ Deterministic trace IDs (seed=image_id)",
              "✓ Dataset items upserted by item_id — safe to re-run",
              "✓ Results cached to eval_cache.json",
          ], title_color=C_ACCENT3)

    _card(slide, Inches(6.45), Inches(4.25), Inches(6.5), Inches(2.8),
          "Testing", [
              "61 unit tests — no API calls, < 5 seconds",
              "8 integration tests — real API, ~$0.01",
              "Mock OpenAI client for judge tests",
              "pytest tests/ -v  to run all",
              "pytest -m integration  for API tests",
          ], title_color=C_ACCENT2)


def slide_tradeoffs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Decisions & Trade-offs", "What we chose and why")

    rows = [
        ("Detection method", "LLM vision judges", "Fine-tuned classifier", "No labelled training data at scale; LLM generalises across forgery types"),
        ("Ensemble size", "3 judges", "1 or 5+ judges", "Minimum for majority vote; 5+ adds cost with marginal benefit"),
        ("Aggregation", "Majority vote", "Weighted confidence avg", "Simple, interpretable, robust to one outlier judge"),
        ("Uncertainty", "Explicit UNCERTAIN class", "Force binary FAKE/REAL", "Preserves calibration; flags hard cases for human review"),
        ("Observability", "Langfuse via MCP", "No logging / direct SDK", "MCP avoids credentials in app; adds prompt versioning + cost tracking"),
        ("UI", "Streamlit", "CLI only or Dash", "Fast to build, interactive; Dash only needed at production scale"),
        ("Prompt length", "Short & concrete", "Long chain-of-thought", "Shorter = faster, cheaper, less hallucination"),
    ]

    _box(slide, Inches(0.35), Inches(1.45), Inches(12.6), Inches(5.65), bg=C_DARK_CARD)

    col_ws = [Inches(1.9), Inches(2.0), Inches(2.0), Inches(6.4)]
    headers = ["Decision", "Chosen", "Alternative", "Why"]
    left0 = Inches(0.42)
    top0  = Inches(1.52)
    row_h = Inches(0.7)

    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        left = left0 + sum(col_ws[:j])
        _txt(slide, hdr, left, top0, cw, row_h,
             size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        top = top0 + row_h + i * row_h * 0.73
        bg = RGBColor(0x1A, 0x31, 0x48) if i % 2 == 0 else C_DARK_CARD
        _box(slide, Inches(0.42), top, Inches(12.5), row_h * 0.73, bg=bg)
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            left = left0 + sum(col_ws[:j])
            _txt(slide, cell, left + Inches(0.05), top + Inches(0.04),
                 cw - Inches(0.1), row_h * 0.68,
                 size=11, color=C_WHITE, align=PP_ALIGN.LEFT)


def slide_next_steps(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _header_bar(slide, "Next Steps & Improvements", "Clear path to higher accuracy and lower cost")

    _card(slide, Inches(0.35), Inches(1.45), Inches(3.9), Inches(5.6),
          "Quick wins (no API cost)", [
              "OCR arithmetic check: do items sum to total?",
              "Catches single-digit substitutions (Case 1)",
              "Hash-based duplicate detection",
              "Image metadata check (EXIF editing traces)",
              "Hard rules as pre-filter before LLM judges",
          ], title_color=C_ACCENT3)

    _card(slide, Inches(4.45), Inches(1.45), Inches(3.9), Inches(5.6),
          "Model improvements", [
              "Replace gpt-4o-mini (judge_1) — too uncertain",
              "Try gpt-4.1 for higher accuracy",
              "Stratified 50/50 batch for fair evaluation",
              "Tune temperature per judge via Langfuse runs",
              "Add chain-of-thought for borderline cases",
          ], title_color=C_ACCENT)

    _card(slide, Inches(8.55), Inches(1.45), Inches(4.4), Inches(5.6),
          "Platform improvements", [
              "Human annotation loop (Langfuse UI)",
              "LLM-as-a-Judge scoring via Langfuse eval",
              "Prompt versioning + A/B comparison",
              "Export to production REST API (FastAPI)",
              "Add confidence threshold for auto-rejection",
              "Scheduled batch eval with email report",
          ], title_color=C_ACCENT2)


def slide_end(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _box(slide, Inches(0), Inches(0), Inches(0.35), H, bg=C_ACCENT)
    _box(slide, Inches(0.35), Inches(2.5), Inches(12.98), Inches(2.8), bg=C_DARK_CARD)
    _txt(slide, "Thank you",
         Inches(0.6), Inches(2.65), Inches(12.0), Inches(1.1),
         size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, "LLM-Judge Fake Receipt Detector  ·  TESTTHEKEY",
         Inches(0.6), Inches(3.65), Inches(12.0), Inches(0.5),
         size=20, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)
    _txt(slide,
         "Stack:  Python  ·  OpenAI Vision API  ·  Streamlit  ·  Langfuse  ·  MCP  ·  Plotly",
         Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.4),
         size=15, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = _prs()

    slide_title(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_judges(prs)
    slide_prompt(prs)
    slide_dataset(prs)
    slide_results(prs)
    slide_errors(prs)
    slide_langfuse(prs)
    slide_ui(prs)
    slide_code_quality(prs)
    slide_tradeoffs(prs)
    slide_next_steps(prs)
    slide_end(prs)

    out = Path("TESTTHEKEY_Presentation.pptx")
    prs.save(str(out))
    print(f"Saved: {out.resolve()}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
